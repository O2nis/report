# autofill_word_pptx.py  –  2025-07-17
# pip install streamlit python-docx python-pptx pandas openpyxl
import streamlit as st, pandas as pd
from io import BytesIO
import zipfile
import re
from docx import Document
from docx.oxml.ns import qn as qn_docx
from docx.shared import RGBColor as DocxRGBColor
from docx.oxml import OxmlElement
from pptx import Presentation
from pptx.dml.color import RGBColor as PptxRGBColor
from pptx.util import Pt, Inches
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
from pptx.oxml.ns import qn as qn_pptx

# ───────────────────────── HELPERS ──────────────────────────
def parse_value_and_color(value):
    """Parse value|color syntax, return (value, color) or (value, None)."""
    if isinstance(value, str) and '|' in value:
        parts = value.split('|', 1)
        if len(parts) == 2 and parts[1].lower() in ('green', 'yellow', 'red'):
            return parts[0], parts[1].lower()
    return value, None

def get_color_rgb(color):
    """Return RGB tuple for color name or None if not recognized."""
    color_map = {'green': (0, 255, 0), 'yellow': (255, 255, 0), 'red': (255, 0, 0)}
    return color_map.get(color)

def join_runs(paragraph):
    """Return full text of paragraph without modifying runs."""
    return ''.join(run.text for run in paragraph.runs)

# ───────────────────────── WORD ──────────────────────────
def replace_in_word(doc: Document, placeholders: dict):
    replacement_count = 0

    def replace_in_paragraph(para, placeholders):
        nonlocal replacement_count
        full_text = join_runs(para)
        if not any(f'{{{k}}}' in full_text for k in placeholders):
            return 0
        count = 0
        for run in para.runs:
            text = run.text
            for k, v in placeholders.items():
                tok = f'{{{k}}}'
                if tok in text:
                    st.write(f"Replacing {tok} in paragraph")
                    value, color = parse_value_and_color(v)
                    count += text.count(tok)
                    run.text = text.replace(tok, str(value))
                    if color and k.startswith('dot|'):
                        run.font.color.rgb = DocxRGBColor(*get_color_rgb(color))
        return count

    def replace_in_table(tbl, placeholders):
        nonlocal replacement_count
        count = 0
        for row in tbl.rows:
            for cell in row.cells:
                for para in cell.paragraphs:
                    full_text = join_runs(para)
                    if any(f'{{{k}}}' in full_text for k in placeholders):
                        count += replace_in_paragraph(para, placeholders)
                        for k, v in placeholders.items():
                            tok = f'{{{k}}}'
                            if tok in full_text and not k.startswith('dot|'):
                                value, color = parse_value_and_color(v)
                                if color:
                                    cell_xml = cell._element
                                    tcPr = cell_xml.get_or_add_tcPr()
                                    shd = tcPr.find(qn_docx('w:shd'))
                                    if shd is None:
                                        shd = OxmlElement('w:shd')
                                        tcPr.append(shd)
                                    shd.set(qn_docx('w:fill'), f'{get_color_rgb(color)[0]:02X}{get_color_rgb(color)[1]:02X}{get_color_rgb(color)[2]:02X}')
        return count

    # Process headers and footers
    for section in doc.sections:
        for header in (section.header, section.first_page_header, section.even_page_header):
            for para in header.paragraphs:
                replacement_count += replace_in_paragraph(para, placeholders)
            for tbl in header.tables:
                replacement_count += replace_in_table(tbl, placeholders)
        for footer in (section.footer, section.first_page_footer, section.even_page_footer):
            for para in footer.paragraphs:
                replacement_count += replace_in_paragraph(para, placeholders)
            for tbl in footer.tables:
                replacement_count += replace_in_table(tbl, placeholders)

    # Main body
    for para in doc.paragraphs:
        replacement_count += replace_in_paragraph(para, placeholders)
    for tbl in doc.tables:
        replacement_count += replace_in_table(tbl, placeholders)

    return doc, replacement_count

def save_word(doc):
    buf = BytesIO(); doc.save(buf); buf.seek(0); return buf

# ─────────────────────── PPTX HELPERS ───────────────────────
def _strip_table_borders(shape):
    """Remove borders inside a PPT table."""
    if not shape.has_table:
        return
    tbl = shape.table
    for row in tbl.rows:
        for cell in row.cells:
            tcPr = cell._tc.get_or_add_tcPr()
            for tag in ("a:lnL", "a:lnR", "a:lnT", "a:lnB", "a:lnTlToBr", "a:lnBlToTr"):
                ln = tcPr.find(qn_pptx(tag))
                if ln is not None:
                    tcPr.remove(ln)

def _purge_dashed_shapes(shapes):
    """Delete only dashed-line shapes, preserving images and other shapes."""
    doomed = []
    for shp in shapes:
        if shp.shape_type == MSO_SHAPE_TYPE.GROUP:
            _purge_dashed_shapes(shp.shapes)
        elif shp.shape_type == MSO_SHAPE_TYPE.LINE:
            try:
                ln = shp.line
                if ln and ln.dash_style is not None:
                    doomed.append(shp)
            except Exception:
                pass
    for shp in doomed:
        try:
            shp._element.getparent().remove(shp._element)
        except Exception:
            pass

def _process_shape_text(shape, placeholders, replacement_count, location="unknown"):
    """Replace tokens in a shape’s text frame, preserving run formatting."""
    if not shape.has_text_frame:
        return replacement_count
    tf = shape.text_frame
    for para in tf.paragraphs:
        full_text = join_runs(para)
        if not any(f'{{{k}}}' in full_text for k in placeholders):
            continue
        for run in para.runs:
            text = run.text
            for k, v in placeholders.items():
                tok = f'{{{k}}}'
                if tok in text:
                    st.write(f"Replacing {tok} in {location}")
                    value, color = parse_value_and_color(v)
                    replacement_count += text.count(tok)
                    run.text = text.replace(tok, str(value))
                    if color and k.startswith('dot|'):
                        run.font.color.rgb = PptxRGBColor(*get_color_rgb(color))
    return replacement_count

def _process_shapes_collection(shapes, placeholders, slide=None, location="unknown"):
    """Handle text, table borders, and dashed lines in a shapes collection, preserving images."""
    replacement_count = 0
    _purge_dashed_shapes(shapes)

    for shp in shapes:
        if shp.shape_type == MSO_SHAPE_TYPE.GROUP:
            replacement_count += _process_shapes_collection(shp.shapes, placeholders, slide, f"{location} group")
        elif shp.shape_type == MSO_SHAPE_TYPE.TABLE:
            _strip_table_borders(shp)
            tbl = shp.table
            for row in tbl.rows:
                for cell in row.cells:
                    for para in cell.text_frame.paragraphs:
                        full_text = join_runs(para)
                        for k, v in placeholders.items():
                            tok = f'{{{k}}}'
                            if tok in full_text:
                                st.write(f"Replacing {tok} in {location} table")
                                value, color = parse_value_and_color(v)
                                replacement_count += full_text.count(tok)
                                for run in para.runs:
                                    if tok in run.text:
                                        run.text = run.text.replace(tok, str(value))
                                        if color and not k.startswith('dot|'):
                                            cell.fill.solid()
                                            cell.fill.fore_color.rgb = PptxRGBColor(*get_color_rgb(color))
        elif shp.shape_type != MSO_SHAPE_TYPE.PICTURE:  # Skip images
            if slide and any(k.startswith('dot|') for k in placeholders.keys()) and shp.has_text_frame:
                for para in shp.text_frame.paragraphs:
                    full_text = join_runs(para)
                    for k, v in placeholders.items():
                        tok = f'{{{k}}}'
                        if tok in full_text and k.startswith('dot|'):
                            st.write(f"Replacing {tok} in {location} for dot")
                            value, color = parse_value_and_color(v)
                            if color:
                                left, top, width, height = shp.left, shp.top, Inches(0.1), Inches(0.1)
                                dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, width, height)
                                dot.fill.solid()
                                dot.fill.fore_color.rgb = PptxRGBColor(*get_color_rgb(color))
                                dot.line.color.rgb = PptxRGBColor(*get_color_rgb(color))
                                replacement_count += full_text.count(tok)
                                for run in para.runs:
                                    run.text = run.text.replace(tok, '')
                            continue
                        replacement_count = _process_shape_text(shp, placeholders, replacement_count, location)
            else:
                replacement_count = _process_shape_text(shp, placeholders, replacement_count, location)
    return replacement_count

def replace_in_pptx(prs: Presentation, placeholders: dict):
    replacement_count = 0
    for i, master in enumerate(prs.slide_masters):
        replacement_count += _process_shapes_collection(master.shapes, placeholders, location=f"slide master {i}")
    for i, layout in enumerate(prs.slide_layouts):
        replacement_count += _process_shapes_collection(layout.shapes, placeholders, location=f"slide layout {i}")
    for i, slide in enumerate(prs.slides):
        replacement_count += _process_shapes_collection(slide.shapes, placeholders, slide, f"slide {i}")
        if hasattr(slide, 'notes_slide'):
            for shape in slide.notes_slide.shapes:
                replacement_count += _process_shape_text(shape, placeholders, replacement_count, f"slide {i} notes")
    return prs, replacement_count

def save_pptx(prs):
    buf = BytesIO(); prs.save(buf); buf.seek(0); return buf

# ───────────────────────── STREAMLIT UI ──────────────────────────
st.set_page_config(page_title="Auto-Fill Word / PPTX", layout="centered")
st.title("📝 Auto-fill Word or PowerPoint templates")

kind = st.radio("Template type:", ("Word (.docx)", "PowerPoint (.pptx)"), horizontal=True)
tfile = st.file_uploader("Upload template",
                         type=["docx"] if kind.startswith("Word") else ["pptx"])
xfile = st.file_uploader("Upload Excel with keywords & values", type=["xlsx"])

report_mode = st.radio("Report generation mode:", 
                       ["Single report", 
                        "Multiple reports (one per column)", 
                        "Multiple reports (one per sheet)"],
                       horizontal=True)

if tfile and xfile:
    xl = pd.ExcelFile(xfile)
    all_sheet_names = xl.sheet_names
    
    if report_mode == "Multiple reports (one per sheet)":
        df_first = pd.read_excel(xfile, sheet_name=all_sheet_names[0])
        if df_first.empty:
            st.error(f"First sheet '{all_sheet_names[0]}' is empty"); st.stop()
        cols = df_first.columns.tolist()
        
        st.subheader("Column Selection (based on first sheet)")
        with st.container():
            col1, col2 = st.columns(2)
            with col1:
                kw_cols = st.multiselect("Keyword columns (select 1-3)", cols, key="kw_per_sheet", max_selections=3)
            with col2:
                val_cols = st.multiselect("Value columns (select 1-3)", cols, key="val_per_sheet", max_selections=3)
        
        if not kw_cols or not val_cols:
            st.error("Please select at least one keyword and one value column")
            st.stop()
        if len(kw_cols) != len(val_cols):
            st.error("Number of keyword columns must match number of value columns")
            st.stop()
    else:
        sheet_name = st.selectbox("Select sheet", all_sheet_names, index=0)
        df = pd.read_excel(xfile, sheet_name=sheet_name)
        if df.empty:
            st.error(f"Sheet '{sheet_name}' is empty"); st.stop()
        cols = df.columns.tolist()
        kw_col = st.selectbox("Keyword column", cols, key="kw")
        
        if report_mode == "Single report":
            val_col = st.selectbox("Value column", cols, key="val_single")
        else:
            val_col = st.selectbox("First value column", cols, key="val_multi")

    if st.button("Generate file(s)"):
        if report_mode == "Single report":
            keys = df[kw_col].astype(str).tolist()
            vals = df[val_col].astype(str).tolist()
            if len(keys) != len(vals):
                st.error(f"Columns {kw_col} and {val_col} must have same number of rows"); st.stop()

            mapping = dict(zip(keys, vals))

            if kind.startswith("Word"):
                doc = Document(tfile)
                filled, count = replace_in_word(doc, mapping)
                buf = save_word(filled)
                st.download_button("⬇️ Download filled Word",
                                   data=buf, file_name="filled.docx",
                                   mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document")
                st.write(f"Report 1 ({val_col}): Replaced {count} keywords")
            else:
                prs = Presentation(tfile)
                filled, count = replace_in_pptx(prs, mapping)
                buf = save_pptx(filled)
                st.download_button("⬇️ Download filled PowerPoint",
                                   data=buf, file_name="filled.pptx",
                                   mime="application/vnd.openxmlformats-officedocument.presentationml.presentation")
                st.write(f"Report 1 ({val_col}): Replaced {count} keywords")
                
        elif report_mode == "Multiple reports (one per column)":
            val_cols = df.columns[df.columns.get_loc(val_col):].tolist()
            zip_buffer = BytesIO()
            with zipfile.ZipFile(zip_buffer, 'w', zipfile.ZIP_DEFLATED) as zf:
                for i, v_col in enumerate(val_cols, 1):
                    keys = df[kw_col].astype(str).tolist()
                    vals = df[v_col].astype(str).tolist()
                    if len(keys) != len(vals):
                        st.error(f"Columns {kw_col} and {v_col} must have same number of rows"); st.stop()

                    mapping = dict(zip(keys, vals))

                    if kind.startswith("Word"):
                        doc = Document(tfile)
                        filled, count = replace_in_word(doc, mapping)
                        buf = save_word(filled)
                        file_name = f"filled_{i}.docx"
                        zf.writestr(file_name, buf.getvalue())
                        st.write(f"Report {i} ({v_col}): Replaced {count} keywords")
                    else:
                        prs = Presentation(tfile)
                        filled, count = replace_in_pptx(prs, mapping)
                        buf = save_pptx(filled)
                        file_name = f"filled_{i}.pptx"
                        zf.writestr(file_name, buf.getvalue())
                        st.write(f"Report {i} ({v_col}): Replaced {count} keywords")
            zip_buffer.seek(0)
            st.download_button("⬇️ Download all reports (ZIP)",
                               data=zip_buffer,
                               file_name="filled_reports.zip",
                               mime="application/zip")
                               
        else:  # Multiple reports per sheet
            zip_buffer = BytesIO()
            with zipfile.ZipFile(zip_buffer, 'w', zipfile.ZIP_DEFLATED) as zf:
                for sheet_name in all_sheet_names:
                    try:
                        df_sheet = pd.read_excel(xfile, sheet_name=sheet_name)
                        missing_kw = [col for col in kw_cols if col not in df_sheet.columns]
                        missing_val = [col for col in val_cols if col not in df_sheet.columns]
                        
                        if missing_kw or missing_val:
                            st.warning(f"Sheet '{sheet_name}' missing columns: "
                                       f"Keywords: {', '.join(missing_kw)} | "
                                       f"Values: {', '.join(missing_val)}. Skipping.")
                            continue
                            
                        mapping = {}
                        for kw_col, val_col in zip(kw_cols, val_cols):
                            if len(df_sheet[kw_col]) != len(df_sheet[val_col]):
                                st.warning(f"Sheet '{sheet_name}': Columns {kw_col} and {val_col} have different row counts. Skipping pair.")
                                continue
                            for key, value in zip(df_sheet[kw_col].astype(str), df_sheet[val_col].astype(str)):
                                if key and key.strip():
                                    mapping[key] = value
                        
                        if not mapping:
                            st.warning(f"Sheet '{sheet_name}': No valid key-value pairs found. Skipping.")
                            continue
                            
                        safe_name = re.sub(r'[^\w\s-]', '', sheet_name).strip().replace(' ', '_')[:50]
                        
                        if kind.startswith("Word"):
                            doc = Document(tfile)
                            filled, count = replace_in_word(doc, mapping)
                            buf = save_word(filled)
                            file_name = f"filled_{safe_name}.docx"
                            zf.writestr(file_name, buf.getvalue())
                            st.write(f"Sheet '{sheet_name}': Replaced {count} keywords")
                        else:
                            prs = Presentation(tfile)
                            filled, count = replace_in_pptx(prs, mapping)
                            buf = save_pptx(filled)
                            file_name = f"filled_{safe_name}.pptx"
                            zf.writestr(file_name, buf.getvalue())
                            st.write(f"Sheet '{sheet_name}': Replaced {count} keywords")
                            
                    except Exception as e:
                        st.error(f"Error processing sheet '{sheet_name}': {str(e)}")
            
            zip_buffer.seek(0)
            st.download_button("⬇️ Download all sheets (ZIP)",
                               data=zip_buffer,
                               file_name="filled_sheets.zip",
                               mime="application/zip")
else:
    st.info("Upload both template and Excel to begin.")
